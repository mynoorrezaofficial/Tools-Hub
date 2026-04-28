import os
import csv
import zipfile
from PIL import Image
import fitz  # PyMuPDF
from docx import Document
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt


# ──────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────

def _read_text(path):
    """Read any text-like file and return as list of lines."""
    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            return f.readlines()
    except Exception:
        return []


def _lines_from_docx(path):
    doc = Document(path)
    return [p.text for p in doc.paragraphs if p.text.strip()]


def _lines_from_pdf(path):
    lines = []
    doc = fitz.open(path)
    for page in doc:
        lines += page.get_text().splitlines()
    doc.close()
    return lines


def _get_text_lines(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == '.docx':
        return _lines_from_docx(path)
    if ext == '.pdf':
        return _lines_from_pdf(path)
    return _read_text(path)


# ──────────────────────────────────────────────
# IMAGE CONVERTERS
# ──────────────────────────────────────────────

def image_to_pdf(input_paths, output_path):
    try:
        if isinstance(input_paths, str):
            input_paths = [input_paths]
        images = []
        for path in input_paths:
            img = Image.open(path)
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            images.append(img)
        if images:
            images[0].save(output_path, "PDF", resolution=100.0,
                           save_all=True, append_images=images[1:])
            return True, "Converted to PDF."
        return False, "No valid images found."
    except Exception as e:
        return False, str(e)


def image_to_image(source_path, output_path, fmt):
    """Convert image to JPG/JPEG/PNG/WEBP."""
    try:
        img = Image.open(source_path)
        save_fmt = 'JPEG' if fmt in ('jpg', 'jpeg') else fmt.upper()
        if img.mode in ("RGBA", "P") and save_fmt == 'JPEG':
            img = img.convert("RGB")
        img.save(output_path, save_fmt)
        return True, f"Converted to {fmt.upper()}."
    except Exception as e:
        return False, str(e)


def pdf_to_single_image(source_path, output_path, fmt):
    """Extract first page of PDF as an image."""
    try:
        doc = fitz.open(source_path)
        if len(doc) == 0:
            return False, "PDF is empty."
        page = doc.load_page(0)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        pix.save(output_path)
        doc.close()
        return True, f"First PDF page extracted as {fmt.upper()}."
    except Exception as e:
        return False, str(e)


def pdf_to_images_zip(pdf_path, output_zip_path, temp_folder):
    """Convert all PDF pages to PNG images, zipped."""
    try:
        doc = fitz.open(pdf_path)
        img_paths = []
        os.makedirs(temp_folder, exist_ok=True)
        with zipfile.ZipFile(output_zip_path, 'w') as zipf:
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_name = f"page_{page_num + 1}.png"
                img_path = os.path.join(temp_folder, img_name)
                pix.save(img_path)
                zipf.write(img_path, img_name)
                img_paths.append(img_path)
        doc.close()
        for p in img_paths:
            if os.path.exists(p):
                os.remove(p)
        return True, "PDF converted to ZIP of images."
    except Exception as e:
        return False, str(e)


# ──────────────────────────────────────────────
# PDF CONVERTERS
# ──────────────────────────────────────────────

def txt_to_pdf(txt_path, output_pdf_path):
    try:
        c = canvas.Canvas(output_pdf_path, pagesize=letter)
        width, height = letter
        c.setFont("Helvetica", 12)
        lines = _read_text(txt_path)
        y = height - 40
        for line in lines:
            if y < 40:
                c.showPage()
                c.setFont("Helvetica", 12)
                y = height - 40
            c.drawString(40, y, line.rstrip()[:120])
            y -= 16
        c.save()
        return True, "Text converted to PDF."
    except Exception as e:
        return False, str(e)


def lines_to_pdf(lines, output_pdf_path):
    """Generic lines → PDF helper."""
    try:
        c = canvas.Canvas(output_pdf_path, pagesize=letter)
        width, height = letter
        c.setFont("Helvetica", 12)
        y = height - 40
        for line in lines:
            if y < 40:
                c.showPage()
                c.setFont("Helvetica", 12)
                y = height - 40
            c.drawString(40, y, str(line).rstrip()[:120])
            y -= 16
        c.save()
        return True, "Converted to PDF."
    except Exception as e:
        return False, str(e)


# ──────────────────────────────────────────────
# DOCX CONVERTERS
# ──────────────────────────────────────────────

def lines_to_docx(lines, output_path, title="Converted Document"):
    try:
        doc = Document()
        doc.add_heading(title, 0)
        for line in lines:
            doc.add_paragraph(str(line).strip())
        doc.save(output_path)
        return True, "Converted to DOCX."
    except Exception as e:
        return False, str(e)


def docx_to_txt_raw(docx_path, output_path):
    try:
        lines = _lines_from_docx(docx_path)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))
        return True, "DOCX converted to TXT."
    except Exception as e:
        return False, str(e)


# ──────────────────────────────────────────────
# XLSX CONVERTERS
# ──────────────────────────────────────────────

def lines_to_xlsx(lines, output_path):
    try:
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        for row_idx, line in enumerate(lines, start=1):
            # Try splitting by comma for CSV-like data
            parts = line.rstrip('\n').split(',')
            for col_idx, val in enumerate(parts, start=1):
                ws.cell(row=row_idx, column=col_idx, value=val.strip())
        wb.save(output_path)
        return True, "Converted to XLSX."
    except Exception as e:
        return False, str(e)


def csv_to_xlsx(csv_path, output_path):
    try:
        wb = openpyxl.Workbook()
        ws = wb.active
        with open(csv_path, newline='', encoding='utf-8', errors='replace') as f:
            reader = csv.reader(f)
            for row in reader:
                ws.append(row)
        wb.save(output_path)
        return True, "CSV converted to XLSX."
    except Exception as e:
        return False, str(e)


# ──────────────────────────────────────────────
# PPTX CONVERTERS
# ──────────────────────────────────────────────

def lines_to_pptx(lines, output_path, title="Converted Presentation"):
    try:
        prs = Presentation()
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if slide.placeholders[1]:
            slide.placeholders[1].text = "Generated by Tools Hub"

        # Content slides — 6 lines per slide
        chunk_size = 6
        for i in range(0, len(lines), chunk_size):
            chunk = lines[i:i + chunk_size]
            content_layout = prs.slide_layouts[1]
            sl = prs.slides.add_slide(content_layout)
            sl.shapes.title.text = f"Section {i // chunk_size + 1}"
            tf = sl.placeholders[1].text_frame
            tf.text = chunk[0] if chunk else ""
            for line in chunk[1:]:
                p = tf.add_paragraph()
                p.text = str(line).strip()
                p.level = 1

        prs.save(output_path)
        return True, "Converted to PPTX."
    except Exception as e:
        return False, str(e)


# ──────────────────────────────────────────────
# TXT CONVERTERS
# ──────────────────────────────────────────────

def lines_to_txt(lines, output_path):
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(str(l).strip() for l in lines))
        return True, "Converted to TXT."
    except Exception as e:
        return False, str(e)


# ──────────────────────────────────────────────
# HTML CONVERTERS
# ──────────────────────────────────────────────

def lines_to_html(lines, output_path, title="Converted Document"):
    try:
        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <title>{title}</title>
  <style>
    body {{ font-family: Inter, sans-serif; max-width: 800px; margin: 40px auto; padding: 0 20px; line-height: 1.6; color: #0f172a; background: #f8fafc; }}
    h1 {{ color: #2563eb; }}
    p {{ background: white; padding: 12px 16px; border-radius: 8px; border: 1px solid #e2e8f0; margin: 8px 0; }}
  </style>
</head>
<body>
  <h1>{title}</h1>
"""
        for line in lines:
            if line.strip():
                html += f"  <p>{line.strip()}</p>\n"
        html += "</body>\n</html>"
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html)
        return True, "Converted to HTML."
    except Exception as e:
        return False, str(e)


# ──────────────────────────────────────────────
# CSV CONVERTERS
# ──────────────────────────────────────────────

def lines_to_csv(lines, output_path):
    try:
        with open(output_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            for line in lines:
                writer.writerow([line.strip()])
        return True, "Converted to CSV."
    except Exception as e:
        return False, str(e)


def xlsx_to_csv(xlsx_path, output_path):
    try:
        wb = openpyxl.load_workbook(xlsx_path)
        ws = wb.active
        with open(output_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            for row in ws.iter_rows(values_only=True):
                writer.writerow([str(v) if v is not None else '' for v in row])
        return True, "XLSX converted to CSV."
    except Exception as e:
        return False, str(e)


# ──────────────────────────────────────────────
# MAIN DISPATCHER
# ──────────────────────────────────────────────

def convert_to_format(input_paths, output_path, target_format, temp_folder):
    """
    Central dispatcher: routes input files → desired output format.
    """
    if not input_paths:
        return False, "No files uploaded."

    source_path = input_paths[0] if isinstance(input_paths, list) else input_paths
    ext = os.path.splitext(source_path)[1].lower().lstrip('.')
    os.makedirs(temp_folder, exist_ok=True)

    IMAGE_EXTS = ('jpg', 'jpeg', 'png', 'webp', 'bmp', 'gif', 'tiff')

    # ── IMAGE OUTPUTS ──────────────────────────────
    if target_format in ('jpg', 'jpeg', 'png', 'webp'):
        if ext == 'pdf':
            return pdf_to_single_image(source_path, output_path, target_format)
        elif ext in IMAGE_EXTS:
            return image_to_image(source_path, output_path, target_format)
        else:
            return False, f"Cannot convert .{ext} to {target_format.upper()}."

    # ── ZIP OUTPUT ────────────────────────────────
    if target_format == 'zip':
        if ext == 'pdf':
            return pdf_to_images_zip(source_path, output_path, temp_folder)
        return False, "ZIP export only supports PDF input."

    # ── PDF OUTPUT ────────────────────────────────
    if target_format == 'pdf':
        if ext in IMAGE_EXTS:
            return image_to_pdf(input_paths, output_path)
        if ext == 'txt':
            return txt_to_pdf(source_path, output_path)
        if ext == 'docx':
            return lines_to_pdf(_lines_from_docx(source_path), output_path)
        if ext in ('csv', 'html'):
            lines = _read_text(source_path)
            return lines_to_pdf(lines, output_path)
        return False, f"Cannot convert .{ext} to PDF."

    # ── DOCX OUTPUT ───────────────────────────────
    if target_format == 'docx':
        if ext in IMAGE_EXTS:
            lines = [f"Image File: {os.path.basename(source_path)}",
                     "Source was an image — text cannot be extracted.",
                     "Generated by Tools Hub."]
        else:
            lines = _get_text_lines(source_path)
        return lines_to_docx(lines, output_path)

    # ── XLSX OUTPUT ───────────────────────────────
    if target_format == 'xlsx':
        if ext == 'csv':
            return csv_to_xlsx(source_path, output_path)
        if ext in IMAGE_EXTS:
            lines = [f"Image File,{os.path.basename(source_path)}",
                     "Note,Generated by Tools Hub"]
        else:
            lines = _get_text_lines(source_path)
        return lines_to_xlsx(lines, output_path)

    # ── PPTX OUTPUT ───────────────────────────────
    if target_format == 'pptx':
        if ext in IMAGE_EXTS:
            lines = [f"Image: {os.path.basename(source_path)}",
                     "Generated by Tools Hub.",
                     "Source was an image file."]
        else:
            lines = _get_text_lines(source_path)
        base = os.path.splitext(os.path.basename(source_path))[0]
        return lines_to_pptx(lines, output_path, title=base)

    # ── TXT OUTPUT ────────────────────────────────
    if target_format == 'txt':
        if ext in IMAGE_EXTS:
            return False, "Cannot extract text from an image file."
        if ext == 'docx':
            return docx_to_txt_raw(source_path, output_path)
        if ext in ('csv', 'html', 'txt'):
            with open(source_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(content)
            return True, "Converted to TXT."
        if ext == 'pdf':
            lines = _lines_from_pdf(source_path)
            return lines_to_txt(lines, output_path)
        return False, f"Cannot convert .{ext} to TXT."

    # ── HTML OUTPUT ───────────────────────────────
    if target_format == 'html':
        if ext in IMAGE_EXTS:
            lines = [f"Image: {os.path.basename(source_path)}",
                     "Generated by Tools Hub."]
        else:
            lines = _get_text_lines(source_path)
        base = os.path.splitext(os.path.basename(source_path))[0]
        return lines_to_html(lines, output_path, title=base)

    # ── CSV OUTPUT ────────────────────────────────
    if target_format == 'csv':
        if ext in IMAGE_EXTS:
            return False, "Cannot extract CSV data from an image file."
        if ext == 'xlsx':
            return xlsx_to_csv(source_path, output_path)
        lines = _get_text_lines(source_path)
        return lines_to_csv(lines, output_path)

    return False, f"Unsupported conversion from .{ext} to {target_format}."
